from pptx import Presentation
from pptx.util import Inches
prs = Presentation()
prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]
notes = [
 ("slide1.png",
  "Gap 1 - the KNOWING gap. 'Uncertain' is the meta-NPS verdict when the five expert "
  "rating systems disagree: 15 of 20 carbohydrate foods. The public then misses on more "
  "than 8 in 10 foods. The miss leans pessimistic (foods read as less healthy than experts "
  "rate them), but about 70% of foods are actually misperceived in BOTH directions - honey "
  "and a sugary 'vitamin C' drink, for instance, read as healthier. Caveat: a conservative "
  "floor - five of many published rating systems, one set of cutoffs."),
 ("slide2.png",
  "Gap 2 - the DOING gap. Fruit-guideline adherence falls from ~61% in toddlers to ~6% in "
  "teens (NHANES 2017-18, NCI usual-intake method): a 55-point cliff the single national "
  "figure (~23% of children) averages away. The income gradient is only 18% to 25%, about "
  "12 cents on the dollar of the age effect; the highest-adherence kids (Mexican-American, "
  "31%) are not the richest. USDA's own report finds health-concern and nutrition-knowledge "
  "behaviors outweigh income and price. Honest limit: this is WHERE the gap lives, not WHY - "
  "fruit intake recovers in adulthood, so it is not a simple autonomy/market story. "
  "Close: both gaps sit in the informational and food environments; willpower and price are "
  "the small dials."),
]
for img, note in notes:
    s = prs.slides.add_slide(blank)
    s.shapes.add_picture(img, 0, 0, width=prs.slide_width, height=prs.slide_height)
    s.notes_slide.notes_text_frame.text = note
prs.save("two-gaps-talk.pptx")
print("saved two-gaps-talk.pptx")
